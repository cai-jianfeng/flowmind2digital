from pptx import Presentation
from pptx.util import Inches
import numpy as np
from sklearn.cluster import KMeans
import math
import random
from pptx.enum.dml import MSO_THEME_COLOR, MSO_LINE
from paddleocr import PaddleOCR
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from predict import predict_mode
from data_preprocess import data_process
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
import sys
from pptx.oxml import parse_xml
from units import geometry

import time


class Canopy:
    def __init__(self, dataset):
        self.dataset = dataset
        self.t1 = 0
        self.t2 = 0

    # 设置初始阈值
    def setThreshold(self, t1, t2):
        if t1 > t2:
            self.t1 = t1
            self.t2 = t2
        else:
            print('t1 needs to be larger than t2!')

    # 使用欧式距离进行距离的计算
    @staticmethod
    def euclideanDistance(vec1, vec2):
        return math.sqrt(((vec1 - vec2)**2).sum())

    # 根据当前dataset的长度随机选择一个下标
    def getRandIndex(self):
        return random.randint(0, len(self.dataset) - 1)

    def clustering(self):
        canopies = []  # 用于存放最终归类结果
        while len(self.dataset) > 1:
            rand_index = self.getRandIndex()
            current_center = self.dataset[rand_index]  # 随机获取一个中心点，定为P点
            current_center_list = []  # 初始化P点的canopy类容器
            delete_list = []  # 初始化P点的删除容器
            self.dataset = np.delete(self.dataset, rand_index,
                                     0)  # 删除随机选择的中心点P
            for datum_j in range(len(self.dataset)):
                datum = self.dataset[datum_j]
                distance = self.euclideanDistance(
                    current_center, datum)  # 计算选取的中心点P到每个点之间的距离
                if distance < self.t1:
                    # 若距离小于t1，则将点归入P点的canopy类
                    current_center_list.append(datum)
                if distance < self.t2:
                    delete_list.append(datum_j)  # 若小于t2则归入删除容器
            # 根据删除容器的下标，将元素从数据集中删除
            self.dataset = np.delete(self.dataset, delete_list, 0)
            canopies.append((current_center, current_center_list))

        k = len(canopies)
        if len(self.dataset) == 1:
            k += 1
        return k

def clustering(X, t1=1.5, t2=0.5, dim=1):
    """
    :param: 一维特征待聚类X 小于t2的分为一类(numpy or list)
    :return:  聚类结果X
    """
    X = np.array(X)
    X = X.reshape(-1, dim)
    gc = Canopy(X)
    gc.setThreshold(t1, t2)
    k = gc.clustering()
    print("t2: ", t2, "k: ", k)
    if k == 1:
        Y = np.zeros(len(X), dtype='int32')
    else:
        Y = KMeans(n_clusters=k).fit_predict(X)
    avg = np.zeros((k, dim))  # 每一类的均值
    cnt = np.zeros((k, dim))  # 每一类的个数
    for x, y in zip(X, Y):
        avg[y] += x
        cnt[y] += 1
    avg = avg / cnt
    ret = np.zeros_like(X)
    for i, y in enumerate(Y):
        ret[i] = avg[y]
    return ret

def align(pred):
    """
    :param pred: boundingbox
    :return:  聚类后的boundingbox
    """
    pred = np.array(pred)  # 长度均值的一半作为t2
    tx = 1e18
    ty = 1e18
    for box in pred:
        tx = min(box[2] - box[0], tx)
        ty = min(box[3] - box[1], ty)
    tx /= 1.618
    ty /= 1.618

    for i in range(4):
        x = pred[:, i]
        if i & 1:
            x = clustering(x, t2=ty)
        else:
            x = clustering(x, t2=tx)
        pred[:, i] = x.reshape(-1)
    return pred.tolist()

def adjust_shape(pred):
    X = np.zeros((len(pred), 2))
    t2 = 1e18
    for i, box in enumerate(pred):
        X[i][0] = box[2] - box[0]
        X[i][1] = box[3] - box[1]
        t2 = min(t2, math.sqrt(X[i][0]**2 + X[i][1]**2))
    X = clustering(X, dim=2, t2=t2/1.618)
    # print(X)

    for i, box in enumerate(pred):
        midx = (box[2] + box[0]) / 2
        box[0] = midx - X[i][0] / 2
        box[2] = midx + X[i][0] / 2

        midy = (box[3] + box[1]) / 2
        box[1] = midy - X[i][1] / 2
        box[3] = midy + X[i][1] / 2
    return pred

def draw_shape(slide, box):
    # 7.33英寸 x 11英寸
    xmin, ymin, xmax, ymax = Inches(box[0]), Inches(box[1]), Inches(box[2]), Inches(box[3])
    category = box[4]

    left = xmin
    top = ymin
    width = xmax - xmin
    height = ymax - ymin
    shapes = slide.shapes

    SHAPE = 0
    COLOR = 0
    """
    'circle': 0,
    'diamonds': 1,
    'long_oval': 2,
    'hexagon': 3,
    'parallelogram': 4,
    'rectangle': 5,
    'trapezoid': 6,
    'triangle': 7,
    """
    if category == 0:  # 椭圆
        SHAPE = MSO_SHAPE.OVAL
        COLOR = MSO_THEME_COLOR.ACCENT_4
    elif category == 1:  # 菱形
        SHAPE = MSO_SHAPE.DIAMOND
        COLOR = MSO_THEME_COLOR.ACCENT_2
    elif category == 2:  # 长椭圆
        SHAPE = MSO_SHAPE.ROUNDED_RECTANGLE
        COLOR = MSO_THEME_COLOR.ACCENT_3
    elif category == 3:  # 六边形
        SHAPE = MSO_SHAPE.HEXAGON
        COLOR = MSO_THEME_COLOR.ACCENT_5
    elif category == 4:  # 平行四边形
        SHAPE = MSO_SHAPE.PARALLELOGRAM
        COLOR = MSO_THEME_COLOR.ACCENT_1
    elif category == 5:  # rectangle
        SHAPE = MSO_SHAPE.RECTANGLE
        COLOR = MSO_THEME_COLOR.ACCENT_1
    elif category == 6:  # 梯形
        SHAPE = MSO_SHAPE.TRAPEZOID
        COLOR = MSO_THEME_COLOR.ACCENT_3
    elif category == 7:
        SHAPE = MSO_SHAPE.ISOSCELES_TRIANGLE
        COLOR = MSO_THEME_COLOR.ACCENT_5




    shape = shapes.add_shape(SHAPE, left, top, width, height)
    fill = shape.fill
    fill.gradient()
    line = shape.line

    fill.gradient_stops[0].color.theme_color = COLOR
    fill.gradient_stops[1].color.theme_color = COLOR
    line.color.theme_color = COLOR
    return shape

def scaling(pred, edge, siz):
    t = max(siz[0] / 7.33, siz[1] / 11)
    for box in pred:
        for i in range(4):
            box[i] /= t
    for e in edge:
        if e[0] == -1:
            e[1] /= t
            e[2] /= t
        if e[3] == -1:
            e[4] /= t
            e[5] /= t
    return pred, edge, t

def work_ocr(scale, img_path):
    ocr = PaddleOCR()
    ocr_result = ocr.ocr(img_path, det=True)  # det + text direction + recog
    ans = ocr_result[0]

    points = list()
    txts = list()
    for item in ans:
        point = item[0]
        p1 = point[0]
        txt = item[1]
        txt = txt[0]
        p1[0] = Inches(p1[0] / scale)
        p1[1] = Inches(p1[1] / scale)
        points.append(p1)
        txts.append(txt)
    return points, txts

def draw(slide, pred, edge, points, txts=None):
    if txts is None:
        txts = list()
    done = np.zeros(len(txts), dtype='int32')
    for box in pred:
        shape = draw_shape(slide, box)
        if not shape.has_text_frame:
            continue
        xmin, ymin, xmax, ymax = Inches(box[0]), Inches(box[1]), Inches(box[2]), Inches(box[3])

        i = 0
        for point, txt in zip(points, txts):
            if xmin - Inches(0.25) <= point[0] <= xmax + Inches(0.25) \
                    and ymin - Inches(0.25) <= point[1] <= ymax + Inches(0.25):
                text_frame = shape.text_frame
                text_frame.text = txt
                # text_frame.margin_bottom = Inches(0.08)
                # text_frame.margin_left = 0
                # text_frame.vertical_anchor = MSO_ANCHOR.TOP
                text_frame.word_wrap = False
                text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                done[i] = 1
                break
            i += 1

    i = -1
    for point, txt in zip(points, txts):
        i += 1
        if done[i] == 1:
            continue
        left = point[0]
        top = point[1]
        txbox = slide.shapes.add_textbox(left, top, Inches(2), Inches(0.5))
        tf = txbox.text_frame
        tf.text = txt

    for e in edge:
        connector_type = MSO_CONNECTOR_TYPE.ELBOW  # 双拐弯折线，可改中间黄点
        if e[0] == -1 or e[3] == -1:
            connector_type = MSO_CONNECTOR_TYPE.STRAIGHT
        connector = slide.shapes.add_connector(
            connector_type, Inches(e[1]), Inches(e[2]), Inches(e[4]), Inches(e[5])
        )
        if e[0] != -1:
            connector.begin_connect(slide.shapes[e[1]], e[2])
        if e[3] != -1:
            connector.end_connect(slide.shapes[e[4]], e[5])
        line = connector.line
        line.dash_style = MSO_LINE.SOLID
        line.color.theme_color = MSO_THEME_COLOR.ACCENT_6
        if e[6] == 9:
            line_elem = line._get_or_add_ln()
            line_elem.append(parse_xml("""
                <a:tailEnd type="arrow" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>
            """))

def model(img_path, opt=0):
    """
    :param img_path:
    :param opt: 0 predict / 1 label
    :return: bbox:(num * 4), np
             cls:(num * 1), np
             keypoint:(num * 2 * 2), np
    """
    predict_util = predict_mode()
    arrow_metadata = predict_util.dataset_register(dataset_path=['./dataset_arrow', './dataset_arrow'])
    predict_output = predict_util.predict_flowchart(arrow_metadata=arrow_metadata,
                                                    img_path=img_path,
                                                    save_path='./output/eval.jpg')
    predict_output = predict_output['instances']
    if opt == 0:
        bbox = predict_output.pred_boxes.tensor
        cls = predict_output.pred_classes
        kpt = predict_output.pred_keypoints
    else:  # 预测
        data_util = data_process()
        bbox, cls, kpt = data_util.get_arrow_dicts(
            dataset_annotation_file='./dataset_arrow/annotations/1212_0355.xml')

    bbox = bbox.cpu().numpy().astype('int32')
    cls = cls.cpu().numpy().astype('int32')
    kpt = kpt.cpu().numpy().astype('int32')
    kpt = kpt[:, :, :-1]
    siz = predict_output.image_size
    return bbox, cls, kpt, siz

def get_pred(bbox, cls):
    ret = list()
    for i, x in enumerate(cls):
        if x >= 8:
            continue
        tmp = bbox[i].tolist()
        tmp.append(x)
        ret.append(tmp)
    return ret

def get_edge(kpt, cls):
    ret = list()
    for i, x in enumerate(cls):
        """
        'arrow': 9,
        'double_arrow': 10,
        'line': 11
        """
        if x < 9:
            continue
        tmp = np.reshape(kpt[i], -1)
        tmp = tmp.tolist()
        tmp.append(x)
        ret.append(tmp)
    return ret

def find_closest_shape(pred, x, y):
    """
    :param pred: all autoshape[xmin, ymin, xmax, ymax, cls]
    :param x:
    :param y:
    :return: opt, id, direction
    """
    mxdis = 1e18
    op = -1
    sp = -1
    d = -1
    for i, shape in enumerate(pred):
        dis, direction = geometry.calc(x, y, shape)  # 计, y到shape最近距离的关键点编号
        if dis < mxdis:
            mxdis = dis
            op = 1   # """之后可以设定最近距离阈值"""
            sp = i
            d = direction
    return op, sp, d

def build_graph(pred, edge):
    """
    :param pred: [[x, y, x, y, cls], ...]
    :param edge: [[x, y, x, y, cls], ...]
    :return: an edge graph[op1, id1, 0-3direction, op2, id2, 0-3direction, edge_cls] 7维
             if opt = -1, use original [x, y] instead of id direction
    """
    ret = list()
    for e in edge:
        op1, sp1, d1 = find_closest_shape(pred, e[0], e[1])
        if op1 == -1:
            sp1, d1 = e[0], e[1]
        op2, sp2, d2 = find_closest_shape(pred, e[2], e[3])
        if op2 == -1:
            sp2, d2 = e[2], e[3]
        cur = [op1, sp1, d1,
               op2, sp2, d2,
               e[4]]
        ret.append(cur)
    return ret


if __name__ == "__main__":
    T0 = time.time()
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # 板式6为空白
    slide = prs.slides.add_slide(slide_layout)

    img_path = './test.jpg'
    bbox, cls, kpt, siz = model(img_path=img_path, opt=0)  # 预测0/ 数据集1

    pred = get_pred(bbox, cls)  # 获得autoshape的框与类别
    edge = get_edge(kpt, cls)  # 获得边的框与类别
    if len(pred) == 0:
        print("not find the shape!")
        sys.exit()

    T1 = time.time()
    print("predict time:", T1 - T0)
    edge = build_graph(pred, edge)  # 建边
    pred, edge, scale = scaling(pred, edge, siz)  # 尺度放缩

    # print("edge:", edge)
    pred = align(pred)  # 对齐聚类排版
    pred = adjust_shape(pred)  # 形状相似聚类排版
    T2 = time.time()
    print("autotypint time:", T2 - T1)
    points, txts = work_ocr(scale, img_path)  # 文字识别 返回点坐标和文本
    T3 = time.time()
    print("OCR time:", T3 - T2)
    print("txts:", txts)

    draw(slide, pred, edge, points, txts)  # 文本 箭头 图形 合并作画
    prs.save("result.pptx")
    T4 = time.time()
    print("generation time:", T4 - T3)
    print("Generated successfully !!!")
